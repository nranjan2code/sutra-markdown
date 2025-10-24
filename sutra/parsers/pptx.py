"""
PPTX Parser - Extract content from PowerPoint presentations

Uses python-pptx for parsing Microsoft PowerPoint presentations with support for:
- Slide text extraction
- Title and content extraction
- Notes extraction
- Table extraction
- Image extraction
- Metadata extraction
- Slide layouts and structure

Requirements:
    pip install python-pptx

Usage:
    parser = PPTXParser("presentation.pptx")
    result = await parser.parse()
    
    # Or streaming
    async for element in parser.parse_stream():
        process(element)
"""

import time
from pathlib import Path
from typing import List, Optional, Dict, Any, AsyncIterator
from datetime import datetime
import asyncio

from ..models.document import (
    ParsedDocument,
    PageInfo,
    ImageInfo,
    TableInfo,
    DocumentElement,
    DocumentMetadata
)
from ..models.enums import ElementType, DocumentType
from .base import BaseParser, ParserResult, ParserError, CorruptedFileError

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False


class PPTXParser(BaseParser):
    """
    PowerPoint presentation parser using python-pptx
    
    Extracts slides, text, tables, images, and metadata from PPTX files.
    Supports both full presentation parsing and streaming mode.
    
    Args:
        file_path: Path to PPTX file
        options: Parser options
            - extract_images: Extract images (default: True)
            - extract_tables: Extract tables (default: True)
            - extract_notes: Extract speaker notes (default: True)
            - include_master_slides: Include master slide info (default: False)
    
    Example:
        >>> parser = PPTXParser("presentation.pptx", {
        ...     "extract_images": True,
        ...     "extract_tables": True,
        ...     "extract_notes": True
        ... })
        >>> result = await parser.parse()
        >>> print(f"Extracted {len(result.document.pages)} slides")
    """
    
    def __init__(
        self,
        file_path: str | Path,
        options: Optional[Dict[str, Any]] = None
    ):
        if not PYTHON_PPTX_AVAILABLE:
            raise ImportError(
                "python-pptx is required for PPTX parsing.\n"
                "Install with: pip install python-pptx"
            )
        
        super().__init__(file_path, options)
        
        # Parse options
        self.extract_images = self.options.get("extract_images", True)
        self.extract_tables = self.options.get("extract_tables", True)
        self.extract_notes = self.options.get("extract_notes", True)
        self.include_master_slides = self.options.get("include_master_slides", False)
        
        # Open PPTX presentation
        try:
            self.prs = Presentation(str(self.file_path))
        except Exception as e:
            raise CorruptedFileError(f"Failed to open PPTX: {e}")
    
    def get_page_count(self) -> int:
        """Get number of slides in presentation"""
        return len(self.prs.slides)
    
    async def validate(self) -> tuple[bool, Optional[str]]:
        """
        Validate PPTX file
        
        Returns:
            Tuple of (is_valid, error_message)
        """
        try:
            # Check if presentation can be opened
            if not hasattr(self, 'prs'):
                return False, "Failed to open PPTX presentation"
            
            # Check if has slides
            if len(self.prs.slides) == 0:
                return False, "PPTX has no slides"
            
            # Try to read first slide
            try:
                slide = self.prs.slides[0]
                _ = slide.shapes
            except Exception as e:
                return False, f"Failed to read PPTX content: {e}"
            
            return True, None
            
        except Exception as e:
            return False, str(e)
    
    async def parse(self) -> ParserResult:
        """
        Parse entire PPTX presentation
        
        Returns:
            ParserResult with parsed content
        """
        start_time = time.time()
        warnings = []
        
        try:
            # Validate first
            is_valid, error = await self.validate()
            if not is_valid:
                return ParserResult(
                    document=None,
                    success=False,
                    error=error,
                    file_size=self.file_size
                )
            
            # Extract metadata
            metadata = self._extract_metadata()
            
            # Parse slides
            pages = []
            images = []
            tables = []
            full_text = []
            
            for slide_idx, slide in enumerate(self.prs.slides):
                slide_data = await self._parse_slide(slide, slide_idx)
                
                # Add page
                pages.append(slide_data["page_info"])
                full_text.append(slide_data["text"])
                
                # Collect images and tables
                images.extend(slide_data["images"])
                tables.extend(slide_data["tables"])
            
            # Create ParsedDocument
            document = ParsedDocument(
                metadata=metadata,
                pages=pages,
                images=images,
                tables=tables,
                full_text="\n\n---\n\n".join(full_text)  # Separate slides with ---
            )
            
            parse_time = time.time() - start_time
            
            return ParserResult(
                document=document,
                success=True,
                warnings=warnings,
                parse_time=parse_time,
                file_size=self.file_size
            )
            
        except Exception as e:
            return ParserResult(
                document=None,
                success=False,
                error=str(e),
                parse_time=time.time() - start_time,
                file_size=self.file_size
            )
    
    async def parse_stream(self) -> AsyncIterator[DocumentElement]:
        """
        Parse PPTX in streaming mode
        
        Yields document elements as they are parsed.
        """
        for slide_idx, slide in enumerate(self.prs.slides):
            slide_data = await self._parse_slide(slide, slide_idx)
            
            # Yield slide text
            if slide_data["text"]:
                yield DocumentElement(
                    type=ElementType.TEXT,
                    content=slide_data["text"],
                    metadata={
                        "slide_number": slide_idx + 1,
                        "title": slide_data.get("title"),
                        "has_notes": slide_data.get("has_notes", False)
                    }
                )
            
            # Yield images
            for image in slide_data["images"]:
                yield DocumentElement(
                    type=ElementType.IMAGE,
                    content=None,
                    metadata={
                        "slide_number": slide_idx + 1,
                        "image_info": image
                    }
                )
            
            # Yield tables
            for table in slide_data["tables"]:
                yield DocumentElement(
                    type=ElementType.TABLE,
                    content=None,
                    metadata={
                        "slide_number": slide_idx + 1,
                        "table_info": table
                    }
                )
            
            # Yield control to event loop
            await asyncio.sleep(0)
    
    async def _parse_slide(self, slide: Any, slide_idx: int) -> Dict[str, Any]:
        """
        Parse a single slide
        
        Args:
            slide: python-pptx Slide object
            slide_idx: Slide index (0-based)
        
        Returns:
            Dict with page_info, text, images, and tables
        """
        text_parts = []
        slide_title = ""
        
        # Extract text from shapes
        for shape in slide.shapes:
            # Check if shape has text
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                
                # Try to identify title
                if shape.is_placeholder and hasattr(shape, "placeholder_format"):
                    try:
                        if shape.placeholder_format.type == 1:  # Title placeholder
                            slide_title = text
                    except:
                        pass
                
                text_parts.append(text)
        
        # Extract notes
        notes_text = ""
        if self.extract_notes and slide.has_notes_slide:
            try:
                notes_slide = slide.notes_slide
                notes_text = notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    text_parts.append(f"\n[Speaker Notes]\n{notes_text}")
            except Exception as e:
                print(f"Warning: Failed to extract notes from slide {slide_idx + 1}: {e}")
        
        # Combine text
        full_slide_text = "\n\n".join(text_parts)
        
        # Extract images
        images = []
        if self.extract_images:
            images = self._extract_images_from_slide(slide, slide_idx)
        
        # Extract tables
        tables = []
        if self.extract_tables:
            tables = self._extract_tables_from_slide(slide, slide_idx)
        
        # Create PageInfo (each slide is a page)
        page_info = PageInfo(
            page_number=slide_idx + 1,
            text=full_slide_text,
            bounds={
                "width": self.prs.slide_width,
                "height": self.prs.slide_height
            }
        )
        
        return {
            "page_info": page_info,
            "text": full_slide_text,
            "title": slide_title,
            "images": images,
            "tables": tables,
            "has_notes": bool(notes_text)
        }
    
    def _extract_images_from_slide(
        self,
        slide: Any,
        slide_idx: int
    ) -> List[ImageInfo]:
        """
        Extract images from slide
        
        Args:
            slide: python-pptx Slide object
            slide_idx: Slide index (0-based)
        
        Returns:
            List of ImageInfo objects
        """
        images = []
        image_index = 0
        
        try:
            for shape in slide.shapes:
                # Check if shape is picture
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        # Get image
                        image = shape.image
                        
                        # Get format
                        ext = image.ext
                        
                        # Get image data
                        image_data = image.blob
                        
                        # Create ImageInfo
                        image_info = ImageInfo(
                            id=f"pptx_image_{slide_idx}_{image_index}",
                            format=ext,
                            width=int(shape.width) if shape.width else 0,
                            height=int(shape.height) if shape.height else 0,
                            page_number=slide_idx + 1,
                            data=image_data
                        )
                        images.append(image_info)
                        image_index += 1
                    
                    except Exception as e:
                        print(f"Warning: Failed to extract image {image_index} from slide {slide_idx + 1}: {e}")
                        continue
        
        except Exception as e:
            print(f"Warning: Failed to get images from slide {slide_idx + 1}: {e}")
        
        return images
    
    def _extract_tables_from_slide(
        self,
        slide: Any,
        slide_idx: int
    ) -> List[TableInfo]:
        """
        Extract tables from slide
        
        Args:
            slide: python-pptx Slide object
            slide_idx: Slide index (0-based)
        
        Returns:
            List of TableInfo objects
        """
        tables = []
        table_index = 0
        
        try:
            for shape in slide.shapes:
                # Check if shape is table
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    try:
                        table = shape.table
                        
                        # Extract rows
                        rows = []
                        for row in table.rows:
                            cells = []
                            for cell in row.cells:
                                cell_text = cell.text.strip()
                                cells.append(cell_text)
                            rows.append(cells)
                        
                        # Create TableInfo (all rows as List[List[str]])
                        table_info = TableInfo(
                            id=f"pptx_table_{slide_idx}_{table_index}",
                            rows=rows,  # All rows including headers
                            page_number=slide_idx + 1,
                            bbox=None,
                            caption=None
                        )
                        tables.append(table_info)
                        table_index += 1
                    
                    except Exception as e:
                        print(f"Warning: Failed to extract table {table_index} from slide {slide_idx + 1}: {e}")
                        continue
        
        except Exception as e:
            print(f"Warning: Failed to get tables from slide {slide_idx + 1}: {e}")
        
        return tables
    
    def _extract_metadata(self) -> DocumentMetadata:
        """
        Extract PPTX metadata
        
        Returns:
            DocumentMetadata object
        """
        try:
            core_props = self.prs.core_properties
            
            # Get dates
            created_date = core_props.created if hasattr(core_props, 'created') else None
            modified_date = core_props.modified if hasattr(core_props, 'modified') else None
            
            # Slide count
            slide_count = len(self.prs.slides)
            
            return self._create_metadata(
                title=core_props.title if hasattr(core_props, 'title') else None,
                author=core_props.author if hasattr(core_props, 'author') else None,
                created_date=created_date,
                modified_date=modified_date,
                page_count=slide_count,
                subject=core_props.subject if hasattr(core_props, 'subject') else None,
                keywords=core_props.keywords if hasattr(core_props, 'keywords') else None
            )
        
        except Exception as e:
            print(f"Warning: Failed to extract metadata: {e}")
            return self._create_metadata(
                page_count=self.get_page_count()
            )
    
    def get_supported_features(self) -> Dict[str, bool]:
        """Get supported features"""
        return {
            "text_extraction": True,
            "table_extraction": self.extract_tables,
            "image_extraction": self.extract_images,
            "metadata_extraction": True,
            "streaming": True,
            "notes": self.extract_notes,
            "master_slides": self.include_master_slides
        }


# Quick test function
async def test_pptx_parser():
    """Test PPTX parser"""
    import sys
    
    if len(sys.argv) < 2:
        print("Usage: python -m sutra.parsers.pptx <pptx_file>")
        return
    
    pptx_file = sys.argv[1]
    
    print(f"\n{'='*70}")
    print(f"🧪 Testing PPTX Parser")
    print(f"{'='*70}")
    print(f"File: {pptx_file}\n")
    
    try:
        # Create parser
        parser = PPTXParser(pptx_file)
        
        print(f"📄 Presentation Info:")
        print(f"   Type: {parser.document_type.value}")
        print(f"   Size: {parser.file_size:,} bytes")
        print(f"   Slides: {parser.get_page_count()}")
        print(f"   Hash: {parser.file_hash[:16]}...")
        
        # Validate
        print(f"\n✓ Validating...")
        is_valid, error = await parser.validate()
        
        if not is_valid:
            print(f"❌ Validation failed: {error}")
            return
        
        print(f"✅ Valid PPTX")
        
        # Parse
        print(f"\n⏳ Parsing...")
        result = await parser.parse()
        
        if not result.success:
            print(f"❌ Parse failed: {result.error}")
            return
        
        print(f"✅ Parsed in {result.parse_time:.2f}s")
        
        # Show results
        doc = result.document
        print(f"\n📊 Results:")
        print(f"   Slides: {len(doc.pages)}")
        print(f"   Images: {len(doc.images)}")
        print(f"   Tables: {len(doc.tables)}")
        print(f"   Total text: {len(doc.full_text):,} chars")
        
        if doc.metadata.title:
            print(f"   Title: {doc.metadata.title}")
        if doc.metadata.author:
            print(f"   Author: {doc.metadata.author}")
        
        # Show first slide preview
        if doc.pages:
            first_slide = doc.pages[0]
            preview = first_slide.text[:200].replace("\n", " ")
            print(f"\n📄 First slide preview:")
            print(f"   {preview}...")
        
        # Show slide details
        if doc.pages:
            print(f"\n📊 Slide breakdown:")
            for idx, page in enumerate(doc.pages[:5]):  # First 5 slides
                word_count = len(page.text.split())
                print(f"   Slide {idx + 1}: {word_count} words")
        
        print(f"\n{'='*70}")
        print(f"✅ Test complete!")
        print(f"{'='*70}\n")
        
    except Exception as e:
        print(f"\n❌ Error: {e}\n")
        import traceback
        traceback.print_exc()


if __name__ == "__main__":
    import asyncio
    asyncio.run(test_pptx_parser())
